# === NON-PROGRAMMER GUIDE ===
# Purpose: Generates PowerPoint and Excel reports from HybridRAG query results.
# What to read first: generate_excel_report() and generate_pptx_report() are
#   the two main entry points. Each takes a list of query results and an output path.
# Inputs: List of QueryResultRecord dicts (query, answer, sources, metrics).
# Outputs: .xlsx or .pptx file saved to the specified path.
# Safety notes: Does not modify any HybridRAG state. Read-only export.
# ============================
# ============================================================================
# HybridRAG -- Report Generator (src/tools/report_generator.py)
# ============================================================================
#
# WHAT THIS FILE DOES (plain English):
#   Takes the answers HybridRAG gave you and turns them into professional
#   PowerPoint slides or Excel spreadsheets you can share with leadership,
#   print for meetings, or attach to project reports.
#
#   Think of it like a secretary: you already did the research (asked
#   questions), now the secretary formats it nicely for your boss.
#
# TWO MAIN FUNCTIONS:
#   generate_excel_report(results, output_path)
#     -> Creates an Excel workbook with:
#        Sheet 1: "Query Results" -- every question, answer, sources, metrics
#        Sheet 2: "Source Summary" -- which documents were cited most
#        Sheet 3: "Report Info" -- when generated, how many queries, settings
#
#   generate_pptx_report(results, output_path)
#     -> Creates a PowerPoint file with:
#        Slide 1: Title slide (report name, date, query count)
#        Slides 2-N: One slide per query result (question + answer + sources)
#        Last slide: Summary statistics
#
# HOW TO USE:
#   From the GUI: Query tab > Export to Excel / Export to PowerPoint
#   From CLI:     python -m src.tools.report_generator --help
#   From code:    from src.tools.report_generator import generate_excel_report
# ============================================================================
from __future__ import annotations

import argparse
import sys
from collections import Counter
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any


@dataclass
class QueryResultRecord:
    """One query + answer pair for report generation.

    This is intentionally separate from QueryResult (query_engine.py) so
    the report generator has no import dependency on the full engine stack.
    The GUI or CLI caller converts QueryResult -> QueryResultRecord before
    calling the generator.
    """
    query: str = ""
    answer: str = ""
    sources: list[dict[str, Any]] = field(default_factory=list)
    chunks_used: int = 0
    tokens_in: int = 0
    tokens_out: int = 0
    cost_usd: float = 0.0
    latency_ms: float = 0.0
    mode: str = "offline"
    timestamp: str = ""
    use_case: str = ""
    error: str = ""


def from_query_result(result, *, query: str = "", use_case: str = "") -> QueryResultRecord:
    """Convert a QueryResult (from query_engine.py) into a QueryResultRecord."""
    return QueryResultRecord(
        query=query,
        answer=getattr(result, "answer", "") or "",
        sources=getattr(result, "sources", []) or [],
        chunks_used=getattr(result, "chunks_used", 0),
        tokens_in=getattr(result, "tokens_in", 0),
        tokens_out=getattr(result, "tokens_out", 0),
        cost_usd=getattr(result, "cost_usd", 0.0),
        latency_ms=getattr(result, "latency_ms", 0.0),
        mode=getattr(result, "mode", "offline") or "offline",
        timestamp=datetime.now().isoformat(),
        use_case=use_case,
        error=getattr(result, "error", "") or "",
    )


# ============================================================================
# Excel Report
# ============================================================================

def generate_excel_report(
    results: list[QueryResultRecord],
    output_path: str | Path,
    *,
    title: str = "HybridRAG Query Report",
) -> Path:
    """Generate a formatted Excel report from query results.

    Returns the Path to the written file.
    """
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    output_path = Path(output_path)

    # -- Colors and styles --
    header_font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="2B579A", end_color="2B579A", fill_type="solid")
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    cell_font = Font(name="Calibri", size=10)
    cell_align = Alignment(vertical="top", wrap_text=True)
    currency_fmt = '#,##0.0000'
    number_fmt = '#,##0'
    thin_border = Border(
        left=Side(style="thin", color="D0D0D0"),
        right=Side(style="thin", color="D0D0D0"),
        top=Side(style="thin", color="D0D0D0"),
        bottom=Side(style="thin", color="D0D0D0"),
    )
    alt_fill = PatternFill(start_color="F2F6FC", end_color="F2F6FC", fill_type="solid")

    # ---- Sheet 1: Query Results ----
    ws = wb.active
    ws.title = "Query Results"
    headers = [
        "Query", "Answer", "Sources", "Chunks", "Tokens In",
        "Tokens Out", "Cost ($)", "Latency (ms)", "Mode", "Use Case",
    ]
    col_widths = [40, 60, 35, 10, 12, 12, 12, 14, 10, 15]

    for col_idx, (header, width) in enumerate(zip(headers, col_widths), 1):
        cell = ws.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
        cell.border = thin_border
        ws.column_dimensions[get_column_letter(col_idx)].width = width

    for row_idx, rec in enumerate(results, 2):
        source_text = _format_sources_text(rec.sources)
        values = [
            rec.query,
            rec.answer,
            source_text,
            rec.chunks_used,
            rec.tokens_in,
            rec.tokens_out,
            rec.cost_usd,
            round(rec.latency_ms),
            rec.mode,
            rec.use_case,
        ]
        row_fill = alt_fill if row_idx % 2 == 0 else None
        for col_idx, value in enumerate(values, 1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.font = cell_font
            cell.alignment = cell_align
            cell.border = thin_border
            if row_fill:
                cell.fill = row_fill
            if col_idx == 7:
                cell.number_format = currency_fmt
            elif col_idx in (4, 5, 6, 8):
                cell.number_format = number_fmt

    ws.auto_filter.ref = "A1:{}{}".format(
        get_column_letter(len(headers)), len(results) + 1)
    ws.freeze_panes = "A2"

    # ---- Sheet 2: Source Summary ----
    ws2 = wb.create_sheet("Source Summary")
    source_counter = Counter()
    for rec in results:
        for src in rec.sources:
            path = src.get("path", "Unknown")
            source_counter[path] += src.get("chunks", 1)

    src_headers = ["Source Document", "Times Cited", "Total Chunks Referenced"]
    src_widths = [50, 15, 22]
    for col_idx, (header, width) in enumerate(zip(src_headers, src_widths), 1):
        cell = ws2.cell(row=1, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_align
        cell.border = thin_border
        ws2.column_dimensions[get_column_letter(col_idx)].width = width

    for row_idx, (path, count) in enumerate(
        source_counter.most_common(), 2
    ):
        cite_count = sum(
            1 for rec in results
            for src in rec.sources
            if src.get("path", "") == path
        )
        for col_idx, value in enumerate([path, cite_count, count], 1):
            cell = ws2.cell(row=row_idx, column=col_idx, value=value)
            cell.font = cell_font
            cell.alignment = cell_align
            cell.border = thin_border
            if row_idx % 2 == 0:
                cell.fill = alt_fill
            if col_idx >= 2:
                cell.number_format = number_fmt

    ws2.auto_filter.ref = "A1:C{}".format(len(source_counter) + 1)
    ws2.freeze_panes = "A2"

    # ---- Sheet 3: Report Info ----
    ws3 = wb.create_sheet("Report Info")
    ws3.column_dimensions["A"].width = 25
    ws3.column_dimensions["B"].width = 40

    info_rows = [
        ("Report Title", title),
        ("Generated", datetime.now().strftime("%Y-%m-%d %H:%M:%S")),
        ("Total Queries", len(results)),
        ("Unique Sources", len(source_counter)),
        ("Total Cost", "${:.4f}".format(sum(r.cost_usd for r in results))),
        ("Avg Latency", "{:.0f} ms".format(
            sum(r.latency_ms for r in results) / max(len(results), 1))),
        ("Online Queries", sum(1 for r in results if r.mode == "online")),
        ("Offline Queries", sum(1 for r in results if r.mode == "offline")),
        ("Errors", sum(1 for r in results if r.error)),
    ]
    for row_idx, (label, value) in enumerate(info_rows, 1):
        cell_a = ws3.cell(row=row_idx, column=1, value=label)
        cell_a.font = Font(name="Calibri", bold=True, size=11)
        cell_b = ws3.cell(row=row_idx, column=2, value=value)
        cell_b.font = cell_font

    wb.save(str(output_path))
    return output_path


# ============================================================================
# PowerPoint Report
# ============================================================================

def generate_pptx_report(
    results: list[QueryResultRecord],
    output_path: str | Path,
    *,
    title: str = "HybridRAG Query Report",
    subtitle: str = "",
) -> Path:
    """Generate a formatted PowerPoint report from query results.

    Returns the Path to the written file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    output_path = Path(output_path)

    # Use widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BLUE = RGBColor(0x2B, 0x57, 0x9A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    DARK = RGBColor(0x33, 0x33, 0x33)
    GRAY = RGBColor(0x66, 0x66, 0x66)
    LIGHT_BG = RGBColor(0xF2, 0xF6, 0xFC)

    # ---- Slide 1: Title ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_filled_rect(slide, 0, 0, prs.slide_width, prs.slide_height, BLUE)

    _add_text_box(
        slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
        title, Pt(36), WHITE, bold=True, alignment=PP_ALIGN.LEFT,
    )
    sub = subtitle or "Generated {}".format(
        datetime.now().strftime("%B %d, %Y at %I:%M %p"))
    _add_text_box(
        slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
        sub, Pt(18), WHITE, alignment=PP_ALIGN.LEFT,
    )
    stats_text = "{} queries | {} unique sources | ${:.4f} total cost".format(
        len(results),
        len(set(
            src.get("path", "")
            for r in results for src in r.sources
        )),
        sum(r.cost_usd for r in results),
    )
    _add_text_box(
        slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
        stats_text, Pt(14), RGBColor(0xBB, 0xCC, 0xEE),
        alignment=PP_ALIGN.LEFT,
    )

    # ---- Slides 2-N: One per query ----
    for idx, rec in enumerate(results, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Header bar
        _add_filled_rect(slide, 0, 0, prs.slide_width, Inches(1.2), BLUE)
        _add_text_box(
            slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
            "Q{}: {}".format(idx, rec.query), Pt(20), WHITE, bold=True,
            alignment=PP_ALIGN.LEFT,
        )

        # Answer box
        answer_text = rec.answer or "(No answer)"
        # Truncate very long answers for readability on slides
        if len(answer_text) > 800:
            answer_text = answer_text[:797] + "..."
        _add_text_box(
            slide, Inches(0.5), Inches(1.5), Inches(8.5), Inches(4.5),
            answer_text, Pt(14), DARK, alignment=PP_ALIGN.LEFT,
        )

        # Sources sidebar
        source_text = "Sources:\n"
        if rec.sources:
            for src in rec.sources[:5]:
                path = src.get("path", "Unknown")
                chunks = src.get("chunks", 0)
                score = src.get("avg_relevance", 0)
                source_text += "  {} ({} chunks, {:.0%})\n".format(
                    Path(path).name, chunks, score)
        else:
            source_text += "  (none)\n"

        _add_text_box(
            slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(3),
            source_text, Pt(11), GRAY, alignment=PP_ALIGN.LEFT,
        )

        # Metrics footer
        metrics = "Mode: {} | Chunks: {} | Latency: {:.0f}ms".format(
            rec.mode, rec.chunks_used, rec.latency_ms)
        if rec.mode == "online":
            metrics += " | Cost: ${:.4f} | Tokens: {}/{}".format(
                rec.cost_usd, rec.tokens_in, rec.tokens_out)
        _add_text_box(
            slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
            metrics, Pt(10), GRAY, alignment=PP_ALIGN.LEFT,
        )

    # ---- Last Slide: Summary ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_filled_rect(slide, 0, 0, prs.slide_width, Inches(1.2), BLUE)
    _add_text_box(
        slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.9),
        "Summary", Pt(28), WHITE, bold=True, alignment=PP_ALIGN.LEFT,
    )

    total_cost = sum(r.cost_usd for r in results)
    avg_latency = sum(r.latency_ms for r in results) / max(len(results), 1)
    online_count = sum(1 for r in results if r.mode == "online")
    offline_count = sum(1 for r in results if r.mode == "offline")
    error_count = sum(1 for r in results if r.error)

    source_counter = Counter()
    for r in results:
        for s in r.sources:
            source_counter[s.get("path", "")] += 1

    summary_lines = [
        "Total queries: {}".format(len(results)),
        "Online: {} | Offline: {}".format(online_count, offline_count),
        "Total cost: ${:.4f}".format(total_cost),
        "Average latency: {:.0f} ms".format(avg_latency),
        "Unique sources cited: {}".format(len(source_counter)),
        "Errors: {}".format(error_count),
    ]
    _add_text_box(
        slide, Inches(0.8), Inches(1.8), Inches(5), Inches(4),
        "\n".join(summary_lines), Pt(16), DARK, alignment=PP_ALIGN.LEFT,
    )

    # Top sources
    if source_counter:
        top_sources = "Most-Cited Documents:\n"
        for path, count in source_counter.most_common(8):
            top_sources += "  {} ({}x)\n".format(Path(path).name, count)
        _add_text_box(
            slide, Inches(6.5), Inches(1.8), Inches(6), Inches(4),
            top_sources, Pt(13), GRAY, alignment=PP_ALIGN.LEFT,
        )

    prs.save(str(output_path))
    return output_path


# ============================================================================
# Helpers
# ============================================================================

def _format_sources_text(sources: list[dict]) -> str:
    if not sources:
        return "(none)"
    parts = []
    for src in sources:
        path = src.get("path", "Unknown")
        chunks = src.get("chunks", 0)
        score = src.get("avg_relevance", 0)
        parts.append("{} ({} chunks, {:.0%})".format(
            Path(path).name, chunks, score))
    return "\n".join(parts)


def _add_text_box(slide, left, top, width, height, text, font_size,
                  color, bold=False, alignment=None):
    from pptx.util import Pt, Emu
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        if alignment:
            p.alignment = alignment
    return txBox


def _add_filled_rect(slide, left, top, width, height, color):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ============================================================================
# CLI entry point
# ============================================================================

def _cli_main():
    parser = argparse.ArgumentParser(
        description="Generate PPT or Excel reports from HybridRAG query results.",
    )
    parser.add_argument(
        "--format", choices=["excel", "pptx", "both"], default="both",
        help="Output format (default: both).",
    )
    parser.add_argument(
        "--output-dir", default="output",
        help="Directory for output files (default: output/).",
    )
    parser.add_argument(
        "--title", default="HybridRAG Query Report",
        help="Report title.",
    )
    parser.add_argument(
        "--queries", nargs="+",
        default=[
            "What is the operating frequency range?",
            "How do leaders and managers differ?",
            "What are the calibration intervals?",
        ],
        help="Queries to run and include in the report.",
    )
    parser.add_argument(
        "--dry-run", action="store_true",
        help="Generate report with sample data instead of running live queries.",
    )
    args = parser.parse_args()

    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    if args.dry_run:
        results = _sample_results()
    else:
        results = _run_queries(args.queries)

    paths = []
    if args.format in ("excel", "both"):
        xlsx_path = output_dir / "hybridrag_report_{}.xlsx".format(timestamp)
        generate_excel_report(results, xlsx_path, title=args.title)
        print("[OK] Excel report: {}".format(xlsx_path))
        paths.append(xlsx_path)

    if args.format in ("pptx", "both"):
        pptx_path = output_dir / "hybridrag_report_{}.pptx".format(timestamp)
        generate_pptx_report(results, pptx_path, title=args.title)
        print("[OK] PowerPoint report: {}".format(pptx_path))
        paths.append(pptx_path)

    return paths


def _sample_results() -> list[QueryResultRecord]:
    """Generate sample data for dry-run testing."""
    return [
        QueryResultRecord(
            query="What is the operating frequency range?",
            answer="The Digisonde-4D operating frequency range is 0.5 - 30 MHz for all modes of operation.",
            sources=[{"path": "Digisonde4DManual.pdf", "chunks": 3, "avg_relevance": 0.89}],
            chunks_used=5, tokens_in=1200, tokens_out=45,
            cost_usd=0.0024, latency_ms=2300, mode="online",
            use_case="Engineering / STEM",
            timestamp=datetime.now().isoformat(),
        ),
        QueryResultRecord(
            query="How do leaders and managers differ?",
            answer="Leaders focus on change and vision while managers focus on stability and processes. Leaders inspire, managers organize.",
            sources=[{"path": "Leadership vs. Management.pdf", "chunks": 2, "avg_relevance": 0.92}],
            chunks_used=4, tokens_in=980, tokens_out=38,
            cost_usd=0.0018, latency_ms=1800, mode="online",
            use_case="Program Management",
            timestamp=datetime.now().isoformat(),
        ),
        QueryResultRecord(
            query="What are the calibration intervals?",
            answer="The maintenance schedule follows quarterly review cycles for calibration intervals.",
            sources=[{"path": "demo_doc.txt", "chunks": 1, "avg_relevance": 0.85}],
            chunks_used=3, tokens_in=0, tokens_out=0,
            cost_usd=0.0, latency_ms=15200, mode="offline",
            use_case="Field Operations",
            timestamp=datetime.now().isoformat(),
        ),
    ]


def _run_queries(queries: list[str]) -> list[QueryResultRecord]:
    """Run live queries through the engine and collect results."""
    project_root = Path(__file__).resolve().parents[2]
    sys.path.insert(0, str(project_root))

    from src.core.bootstrap.boot_coordinator import BootCoordinator
    from src.core.bootstrap.backend_loader import BackendLoader

    bc = BootCoordinator(str(project_root))
    boot_report = bc.run()
    config = boot_report.config

    loader = BackendLoader(config=config, boot_result=boot_report.boot_result)
    bundle = loader.load(timeout_seconds=60)

    if not bundle.query_engine:
        print("[FAIL] No query engine available")
        return []

    results = []
    for query in queries:
        print("  Running: {}".format(query))
        result = bundle.query_engine.query(query)
        rec = from_query_result(result, query=query)
        results.append(rec)
        print("    [OK] {} chunks, {:.0f}ms".format(rec.chunks_used, rec.latency_ms))

    return results


if __name__ == "__main__":
    _cli_main()
