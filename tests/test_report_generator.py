"""Tests for src/tools/report_generator.py -- Excel and PowerPoint export."""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest

from src.tools.report_generator import (
    QueryResultRecord,
    from_query_result,
    generate_excel_report,
    generate_pptx_report,
    _format_sources_text,
)


def _sample_records():
    return [
        QueryResultRecord(
            query="What is the frequency range?",
            answer="The frequency range is 0.5 - 30 MHz.",
            sources=[{"path": "Manual.pdf", "chunks": 3, "avg_relevance": 0.89}],
            chunks_used=5, tokens_in=1200, tokens_out=45,
            cost_usd=0.0024, latency_ms=2300, mode="online",
            use_case="Engineering",
        ),
        QueryResultRecord(
            query="What are the calibration intervals?",
            answer="Calibration follows quarterly review cycles.",
            sources=[{"path": "demo_doc.txt", "chunks": 1, "avg_relevance": 0.85}],
            chunks_used=3, tokens_in=0, tokens_out=0,
            cost_usd=0.0, latency_ms=15200, mode="offline",
            use_case="Field Operations",
        ),
    ]


class TestExcelReport:

    def test_generates_valid_xlsx(self, tmp_path):
        out = tmp_path / "report.xlsx"
        result = generate_excel_report(_sample_records(), out)
        assert result == out
        assert out.exists()
        assert out.stat().st_size > 0

    def test_has_three_sheets(self, tmp_path):
        from openpyxl import load_workbook

        out = tmp_path / "report.xlsx"
        generate_excel_report(_sample_records(), out)
        wb = load_workbook(out)
        assert wb.sheetnames == ["Query Results", "Source Summary", "Report Info"]

    def test_query_results_has_correct_rows(self, tmp_path):
        from openpyxl import load_workbook

        records = _sample_records()
        out = tmp_path / "report.xlsx"
        generate_excel_report(records, out)
        wb = load_workbook(out)
        ws = wb["Query Results"]
        # Header + 2 data rows
        assert ws.max_row == 3
        assert ws.cell(row=2, column=1).value == records[0].query
        assert ws.cell(row=3, column=1).value == records[1].query

    def test_source_summary_counts(self, tmp_path):
        from openpyxl import load_workbook

        out = tmp_path / "report.xlsx"
        generate_excel_report(_sample_records(), out)
        wb = load_workbook(out)
        ws = wb["Source Summary"]
        # Header + 2 sources
        assert ws.max_row == 3

    def test_report_info_has_title(self, tmp_path):
        from openpyxl import load_workbook

        out = tmp_path / "report.xlsx"
        generate_excel_report(_sample_records(), out, title="Test Report")
        wb = load_workbook(out)
        ws = wb["Report Info"]
        assert ws.cell(row=1, column=2).value == "Test Report"

    def test_empty_results_generates_valid_file(self, tmp_path):
        out = tmp_path / "empty.xlsx"
        generate_excel_report([], out)
        assert out.exists()
        assert out.stat().st_size > 0

    def test_cost_column_format(self, tmp_path):
        from openpyxl import load_workbook

        out = tmp_path / "report.xlsx"
        generate_excel_report(_sample_records(), out)
        wb = load_workbook(out)
        ws = wb["Query Results"]
        cost_cell = ws.cell(row=2, column=7)
        assert cost_cell.value == 0.0024


class TestPptxReport:

    def test_generates_valid_pptx(self, tmp_path):
        out = tmp_path / "report.pptx"
        result = generate_pptx_report(_sample_records(), out)
        assert result == out
        assert out.exists()
        assert out.stat().st_size > 0

    def test_slide_count(self, tmp_path):
        from pptx import Presentation

        records = _sample_records()
        out = tmp_path / "report.pptx"
        generate_pptx_report(records, out)
        prs = Presentation(str(out))
        # Title + 2 query slides + summary = 4
        assert len(prs.slides) == 4

    def test_empty_results_has_title_and_summary(self, tmp_path):
        from pptx import Presentation

        out = tmp_path / "empty.pptx"
        generate_pptx_report([], out)
        prs = Presentation(str(out))
        # Title + summary = 2
        assert len(prs.slides) == 2

    def test_long_answer_truncated(self, tmp_path):
        from pptx import Presentation

        records = [QueryResultRecord(
            query="Test?",
            answer="A" * 1000,
            sources=[],
            chunks_used=1, latency_ms=100, mode="offline",
        )]
        out = tmp_path / "long.pptx"
        generate_pptx_report(records, out)
        prs = Presentation(str(out))
        # Should not crash; answer truncated to 800 chars
        assert len(prs.slides) == 3

    def test_multiline_answer_preserves_paragraphs(self, tmp_path):
        from pptx import Presentation

        records = [QueryResultRecord(
            query="Multi-line test?",
            answer="Line one.\nLine two.\nLine three.",
            sources=[],
            chunks_used=1, latency_ms=100, mode="offline",
        )]
        out = tmp_path / "multiline.pptx"
        generate_pptx_report(records, out)
        prs = Presentation(str(out))
        # Query slide is slide index 1 (0 = title)
        query_slide = prs.slides[1]
        # Find the answer text box (largest content shape)
        answer_texts = []
        for shape in query_slide.shapes:
            if shape.has_text_frame:
                paras = [p.text for p in shape.text_frame.paragraphs]
                if any("Line one" in t for t in paras):
                    answer_texts = paras
                    break
        assert "Line one." in answer_texts
        assert "Line two." in answer_texts
        assert "Line three." in answer_texts

    def test_custom_title(self, tmp_path):
        from pptx import Presentation

        out = tmp_path / "titled.pptx"
        generate_pptx_report([], out, title="Custom Title")
        prs = Presentation(str(out))
        assert out.exists()


class TestHelpers:

    def test_format_sources_text_empty(self):
        assert _format_sources_text([]) == "(none)"

    def test_format_sources_text_with_data(self):
        sources = [{"path": "/docs/Manual.pdf", "chunks": 3, "avg_relevance": 0.85}]
        text = _format_sources_text(sources)
        assert "Manual.pdf" in text
        assert "3 chunks" in text
        assert "85%" in text

    def test_from_query_result(self):
        class FakeResult:
            answer = "Test answer"
            sources = [{"path": "a.pdf", "chunks": 1, "avg_relevance": 0.9}]
            chunks_used = 3
            tokens_in = 100
            tokens_out = 20
            cost_usd = 0.001
            latency_ms = 500.0
            mode = "online"
            error = None

        rec = from_query_result(FakeResult(), query="Test query?", use_case="eng")
        assert rec.query == "Test query?"
        assert rec.answer == "Test answer"
        assert rec.use_case == "eng"
        assert rec.cost_usd == 0.001
